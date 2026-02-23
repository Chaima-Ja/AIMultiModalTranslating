"""DOCX and PPTX text extraction."""
from docx import Document
from pptx import Presentation
from typing import List
import sys
import os
import zipfile
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from models import DocxBlock, PptxBlock, ExtractedDocument


def _is_valid_docx(file_path: str) -> bool:
    """Check if file is a valid .docx file (Office Open XML format)."""
    try:
        # .docx files are actually ZIP archives
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            # Check for required files in a valid .docx
            namelist = zip_ref.namelist()
            return 'word/document.xml' in namelist
    except (zipfile.BadZipFile, OSError):
        return False


def extract_docx(file_path: str) -> ExtractedDocument:
    """Extract text blocks from a Word document."""
    # Check if file has .doc extension (old binary format)
    if file_path.lower().endswith('.doc'):
        raise ValueError(
            "The .doc format (old binary Word format) is not supported. "
            "python-docx only supports .docx (Office Open XML) format. "
            "Please convert your file to .docx format first. "
            "You can do this by:\n"
            "1. Opening the file in Microsoft Word and saving as .docx\n"
            "2. Using LibreOffice: soffice --convert-to docx yourfile.doc\n"
            "3. Using online converters"
        )
    
    # Validate that it's actually a .docx file
    if not _is_valid_docx(file_path):
        raise ValueError(
            f"The file '{os.path.basename(file_path)}' is not a valid Word document (.docx). "
            "It may be:\n"
            "- A corrupted file\n"
            "- An old .doc format file (not supported)\n"
            "- A file with incorrect extension\n"
            "Please ensure the file is a valid .docx (Office Open XML) format."
        )
    
    blocks: List[DocxBlock] = []
    try:
        doc = Document(file_path)
    except Exception as e:
        error_msg = str(e)
        if "not a Word file" in error_msg.lower() or "badzipfile" in error_msg.lower():
            raise ValueError(
                f"The file '{os.path.basename(file_path)}' is not a valid Word document. "
                "python-docx only supports .docx (Office Open XML) format, not the old .doc (binary) format. "
                "Please convert your file to .docx format first."
            ) from e
        raise
    
    # Extract paragraphs
    for para_idx, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        
        style_name = para.style.name if para.style else "Normal"
        is_header = style_name.startswith("Heading")
        
        block_id = f"para_{para_idx}"
        
        block = DocxBlock(
            block_id=block_id,
            text=text,
            page=0,  # DOCX doesn't have explicit page numbers
            bbox=(0, 0, 0, 0),  # Not used for DOCX
            font_size=12.0,  # Default, could extract from runs if needed
            font_name="Calibri",  # Default Word font
            is_header=is_header,
            style_name=style_name,
            paragraph_index=para_idx,
            is_table_cell=False,
            table_index=-1,
            row_index=-1,
            col_index=-1
        )
        blocks.append(block)
    
    # Extract table cells
    for tbl_idx, table in enumerate(doc.tables):
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                text = cell.text.strip()
                if not text:
                    continue
                
                block_id = f"tbl{tbl_idx}_r{row_idx}_c{col_idx}"
                
                block = DocxBlock(
                    block_id=block_id,
                    text=text,
                    page=0,
                    bbox=(0, 0, 0, 0),
                    font_size=12.0,
                    font_name="Calibri",
                    is_header=False,  # Table cells are not headers
                    style_name="Normal",
                    paragraph_index=-1,  # Not a paragraph
                    is_table_cell=True,
                    table_index=tbl_idx,
                    row_index=row_idx,
                    col_index=col_idx
                )
                blocks.append(block)
    
    return ExtractedDocument(
        source_path=file_path,
        format="docx",
        blocks=blocks,
        metadata={}
    )


def extract_pptx(file_path: str) -> ExtractedDocument:
    """Extract text blocks from a PowerPoint presentation."""
    blocks: List[PptxBlock] = []
    prs = Presentation(file_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # Skip shapes without text frames
            if not shape.has_text_frame:
                continue
            
            # Get placeholder type if available
            placeholder_type = "NONE"
            if shape.placeholder:
                placeholder_type = str(shape.placeholder.placeholder_format.type)
            
            # Check if this is a title placeholder
            is_title_placeholder = "TITLE" in placeholder_type
            
            # Extract paragraphs from text frame
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                
                # Determine if header: title placeholder or large font
                font_size = 12.0
                if para.runs:
                    font_size = para.runs[0].font.size.pt if para.runs[0].font.size else 12.0
                
                is_header = is_title_placeholder or font_size >= 24
                
                block_id = f"s{slide_idx}_sh{shape_idx}_p{para_idx}"
                
                block = PptxBlock(
                    block_id=block_id,
                    text=text,
                    page=slide_idx,
                    bbox=(0, 0, 0, 0),  # Not used for PPTX
                    font_size=font_size,
                    font_name="Calibri",  # Default
                    is_header=is_header,
                    slide_index=slide_idx,
                    shape_index=shape_idx,
                    shape_name=shape.name,
                    paragraph_index=para_idx,
                    placeholder_type=placeholder_type
                )
                blocks.append(block)
    
    return ExtractedDocument(
        source_path=file_path,
        format="pptx",
        blocks=blocks,
        metadata={}
    )

