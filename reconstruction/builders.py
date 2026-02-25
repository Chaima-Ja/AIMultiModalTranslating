"""Document reconstruction - rebuild files with translated text."""
import os
import shutil
from typing import Dict, Optional
import sys
import numpy as np
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from docx import Document
from pptx import Presentation
import pdfplumber
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
from reportlab.lib.enums import TA_LEFT

from models import ExtractedDocument, DocxBlock, PptxBlock, AudioBlock
from config import Config

# TTS imports (optional - will fail gracefully if not installed)
# Note: Coqui TTS requires Python <3.12. On Python 3.12+, TTS features are disabled.
# If TTS is installed in a separate venv (e.g., venv-tts), run the application
# using that Python interpreter to enable TTS features.
HAS_TTS = False
TTS = None
try:
    from TTS.api import TTS
    HAS_TTS = True
    print("TTS library detected - audio synthesis enabled")
except ImportError:
    # TTS not available - will fall back to SRT generation
    print("TTS library not found - audio files will output as SRT subtitles")
    print("To enable TTS: Install TTS in Python 3.9-3.11 environment or use venv-tts")
except Exception as e:
    # Other errors (e.g., Python version incompatibility, missing dependencies)
    print(f"Warning: TTS not available ({e}). Audio output will use SRT subtitles instead.")
    HAS_TTS = False

# Audio processing imports
try:
    import soundfile as sf
    HAS_SOUNDFILE = True
except ImportError:
    HAS_SOUNDFILE = False

try:
    from scipy.io import wavfile
    HAS_SCIPY = True
except ImportError:
    HAS_SCIPY = False


def _replace_paragraph_text(para, new_text: str):
    """
    Replace paragraph text while preserving first run's formatting.
    
    Sets the first Run's text to the translated string and clears all other runs.
    """
    if not para.runs:
        # No runs exist, add one
        para.add_run(new_text)
        return
    
    # Set first run's text
    para.runs[0].text = new_text
    
    # Clear all other runs
    for run in para.runs[1:]:
        run.text = ""


def rebuild_docx(
    extracted: ExtractedDocument,
    translations: Dict[str, str],
    output_path: str
) -> str:
    """
    Rebuild DOCX file with translated text.
    
    Opens a copy of the original, looks up translations by block_id,
    and replaces paragraph/cell text while preserving formatting.
    """
    # Copy original file
    shutil.copy2(extracted.source_path, output_path)
    
    doc = Document(output_path)
    
    # Create mapping of block_id to translation
    trans_map = translations
    
    # Process paragraphs
    for block in extracted.blocks:
        if not isinstance(block, DocxBlock):
            continue
        
        if block.is_table_cell:
            # Table cell
            if (block.table_index < len(doc.tables) and
                block.row_index < len(doc.tables[block.table_index].rows) and
                block.col_index < len(doc.tables[block.table_index].rows[block.row_index].cells)):
                cell = doc.tables[block.table_index].rows[block.row_index].cells[block.col_index]
                if cell.paragraphs:
                    translated = trans_map.get(block.block_id, block.text)
                    _replace_paragraph_text(cell.paragraphs[0], translated)
        else:
            # Regular paragraph
            if block.paragraph_index < len(doc.paragraphs):
                translated = trans_map.get(block.block_id, block.text)
                _replace_paragraph_text(doc.paragraphs[block.paragraph_index], translated)
    
    doc.save(output_path)
    return output_path


def _replace_pptx_paragraph(para, new_text: str):
    """Replace PPTX paragraph text while preserving first run's formatting."""
    if not para.runs:
        para.add_run().text = new_text
        return
    
    para.runs[0].text = new_text
    for run in para.runs[1:]:
        run.text = ""


def rebuild_pptx(
    extracted: ExtractedDocument,
    translations: Dict[str, str],
    output_path: str
) -> str:
    """
    Rebuild PPTX file with translated text.
    
    Opens a copy of the original, navigates by slide/shape/paragraph indices,
    and replaces text while preserving formatting.
    """
    # Copy original file
    shutil.copy2(extracted.source_path, output_path)
    
    prs = Presentation(output_path)
    trans_map = translations
    
    # Process blocks
    for block in extracted.blocks:
        if not isinstance(block, PptxBlock):
            continue
        
        if (block.slide_index < len(prs.slides) and
            block.shape_index < len(prs.slides[block.slide_index].shapes)):
            shape = prs.slides[block.slide_index].shapes[block.shape_index]
            
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if block.paragraph_index < len(text_frame.paragraphs):
                    translated = trans_map.get(block.block_id, block.text)
                    _replace_pptx_paragraph(text_frame.paragraphs[block.paragraph_index], translated)
    
    prs.save(output_path)
    return output_path


def rebuild_pdf(
    extracted: ExtractedDocument,
    translations: Dict[str, str],
    output_path: str
) -> str:
    """
    Rebuild PDF file with translated text using ReportLab.
    
    Preserves document structure, text wrapping within bounding boxes, and formatting.
    Uses absolute positioning to maintain original layout while allowing text to wrap.
    """
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.colors import black
    from reportlab.lib.enums import TA_LEFT
    from reportlab.platypus import Paragraph
    
    # Read original PDF to get page sizes
    page_sizes = []
    with pdfplumber.open(extracted.source_path) as pdf:
        for page in pdf.pages:
            width = page.width
            height = page.height
            page_sizes.append((width, height))
    
    if not page_sizes:
        # Default to A4 if no pages found
        page_sizes = [(595, 842)]  # A4 in points
    
    # Create new PDF with ReportLab Canvas
    c = canvas.Canvas(output_path, pagesize=page_sizes[0] if page_sizes else letter)
    trans_map = translations
    
    # Group blocks by page
    blocks_by_page = {}
    for block in extracted.blocks:
        page_num = block.page - 1  # Convert 1-indexed to 0-indexed
        if page_num not in blocks_by_page:
            blocks_by_page[page_num] = []
        blocks_by_page[page_num].append(block)
    
    # Sort blocks within each page by Y position (top to bottom), then X (left to right)
    # This ensures we draw them in the correct reading order
    for page_num in blocks_by_page:
        blocks_by_page[page_num].sort(key=lambda b: (-b.bbox[1], b.bbox[0]))  # Sort by Y descending (top first), then X
    
    # Create paragraph styles for text wrapping
    styles = getSampleStyleSheet()
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=11,
        leading=13,
        alignment=TA_LEFT,
        textColor=black,
        wordWrap='CJK',  # Enable word wrapping
    )
    
    header_style = ParagraphStyle(
        'CustomHeader',
        parent=styles['Heading1'],
        fontSize=14,
        leading=16,
        alignment=TA_LEFT,
        textColor=black,
        wordWrap='CJK',
    )
    
    # Draw each page
    for page_num in sorted(blocks_by_page.keys()):
        if page_num > 0:
            c.showPage()
        
        # Set page size if available
        if page_num < len(page_sizes):
            width, height = page_sizes[page_num]
            c.setPageSize((width, height))
        else:
            width, height = page_sizes[0] if page_sizes else (595, 842)
        
        # Get blocks for this page, already sorted
        page_blocks = blocks_by_page[page_num]
        
        # Draw blocks for this page
        for block in page_blocks:
            translated = trans_map.get(block.block_id, block.text)
            
            if not translated.strip():
                continue
            
            # Get bounding box (pdfplumber uses top-left origin)
            x0, y0, x1, y1 = block.bbox
            block_width = x1 - x0
            block_height = y1 - y0
            
            # Convert to ReportLab coordinates (bottom-left origin)
            # pdfplumber: (0,0) is top-left, y increases downward
            # ReportLab: (0,0) is bottom-left, y increases upward
            page_height = height
            rl_x = x0
            rl_y = page_height - y1  # Bottom of the text box
            
            # Determine font size and style
            font_size = max(block.font_size, 10)  # Minimum 10pt
            is_header = block.is_header
            
            # Update styles with block-specific font size
            if is_header:
                style = ParagraphStyle(
                    'BlockHeader',
                    parent=header_style,
                    fontSize=font_size,
                    leading=font_size * 1.2,
                )
            else:
                style = ParagraphStyle(
                    'BlockNormal',
                    parent=normal_style,
                    fontSize=font_size,
                    leading=font_size * 1.2,
                )
            
            # Escape special characters for ReportLab XML
            translated_escaped = translated.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            # Preserve explicit line breaks
            translated_escaped = translated_escaped.replace("\n", "<br/>")
            
            # Create paragraph with text wrapping
            para = Paragraph(translated_escaped, style)
            
            # Get the wrapped height of the paragraph
            # We need to know how much vertical space the text will take
            para_width = block_width
            para_height = para.wrap(para_width, page_height)[1]
            
            # Draw the paragraph at the correct position
            # Use absolute positioning but allow text to wrap within the bounding box
            para.drawOn(c, rl_x, rl_y)
    
    c.save()
    return output_path


def rebuild_audio_output(
    extracted: ExtractedDocument,
    translations: Dict[str, str],
    output_path: str,
    config: Optional[Config] = None,
    generate_audio: bool = True
) -> str:
    """
    Generate translated audio file or SRT subtitle file from translated audio blocks.
    
    If generate_audio=True and TTS is available, creates a synthesized audio file.
    Otherwise, creates an SRT subtitle file.
    
    Args:
        extracted: ExtractedDocument with AudioBlocks
        translations: Dictionary mapping block_id to translated text
        output_path: Output file path
        config: Configuration object (required for TTS)
        generate_audio: If True, generate audio file; if False, generate SRT
    
    Returns:
        Path to generated output file
    """
    trans_map = translations
    
    # Determine output format from file extension
    output_ext = os.path.splitext(output_path)[1].lower()
    if output_ext in ['.mp3', '.wav', '.m4a', '.ogg']:
        generate_audio = True
    elif output_ext == '.srt':
        generate_audio = False
    
    # Generate audio if requested and TTS is available
    if generate_audio and HAS_TTS and config:
        try:
            return _generate_translated_audio(extracted, translations, output_path, config)
        except Exception as e:
            print(f"Warning: TTS audio generation failed ({e}). Falling back to SRT file.")
            # Change extension to .srt
            output_path = os.path.splitext(output_path)[0] + '.srt'
            return _generate_srt_file(extracted, translations, output_path)
    else:
        # Fallback to SRT generation
        if generate_audio and not HAS_TTS:
            print("Warning: TTS not available (requires Python <3.12 or TTS not installed). Generating SRT file instead.")
            # Change extension to .srt
            output_path = os.path.splitext(output_path)[0] + '.srt'
        return _generate_srt_file(extracted, translations, output_path)


def _generate_srt_file(
    extracted: ExtractedDocument,
    translations: Dict[str, str],
    output_path: str
) -> str:
    """Generate SRT subtitle file from translated audio blocks."""
    trans_map = translations
    
    with open(output_path, "w", encoding="utf-8") as f:
        for idx, block in enumerate(extracted.blocks, start=1):
            if not isinstance(block, AudioBlock):
                continue
            
            translated = trans_map.get(block.block_id, block.text)
            
            # Format timestamps for SRT (HH:MM:SS,mmm)
            start_ts = _format_srt_timestamp(block.start_time)
            end_ts = _format_srt_timestamp(block.end_time)
            
            f.write(f"{idx}\n")
            f.write(f"{start_ts} --> {end_ts}\n")
            f.write(f"{translated}\n")
            f.write("\n")
    
    return output_path


def _generate_translated_audio(
    extracted: ExtractedDocument,
    translations: Dict[str, str],
    output_path: str,
    config: Config
) -> str:
    """
    Generate synthesized French audio from translated text blocks.
    
    Uses Coqui TTS to synthesize speech from translated text, preserving
    the original timing structure from Whisper segments.
    """
    if not HAS_TTS:
        raise ImportError(
            "TTS library not installed. Install with: pip install TTS"
        )
    
    if not HAS_SOUNDFILE:
        raise ImportError(
            "soundfile library not installed. Install with: pip install soundfile"
        )
    
    print(f"Loading TTS model: {config.tts_model}")
    try:
        tts = TTS(model_name=config.tts_model, progress_bar=False, gpu=(config.tts_device == "cuda"))
    except Exception as e:
        raise RuntimeError(f"Failed to load TTS model: {e}. Make sure the model is downloaded.") from e
    
    # Collect all translated segments with their timing
    segments = []
    for block in extracted.blocks:
        if not isinstance(block, AudioBlock):
            continue
        
        translated = translations.get(block.block_id, block.text)
        if not translated.strip():
            continue
        
        segments.append({
            'text': translated,
            'start_time': block.start_time,
            'end_time': block.end_time,
            'duration': block.end_time - block.start_time
        })
    
    if not segments:
        raise ValueError("No translated segments found")
    
    # Synthesize audio for each segment
    audio_segments = []
    sample_rate = 22050  # Default TTS sample rate
    
    print(f"Synthesizing {len(segments)} audio segments...")
    for i, seg in enumerate(segments):
        print(f"  Segment {i+1}/{len(segments)}: {seg['text'][:50]}...")
        
        # Generate audio for this segment
        try:
            # Coqui TTS API - XTTS-v2 supports multilingual TTS
            # Use tts_to_file or tts method depending on model
            import tempfile
            
            # Create temporary file for TTS output
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp_file:
                tmp_path = tmp_file.name
            
            try:
                # Try with language parameter (XTTS-v2)
                if hasattr(tts, 'tts_to_file'):
                    tts.tts_to_file(
                        text=seg['text'],
                        file_path=tmp_path,
                        language=config.tts_language
                    )
                elif hasattr(tts, 'tts'):
                    # Some models use tts method directly
                    wav = tts.tts(text=seg['text'], language=config.tts_language)
                    # Save to temp file
                    sf.write(tmp_path, wav, sample_rate)
                else:
                    raise RuntimeError("TTS model does not support expected API")
                
                # Load the generated audio
                wav, sr = sf.read(tmp_path)
                if sr != sample_rate:
                    # Resample if needed (simple approach - just use the sample rate we get)
                    sample_rate = sr
                
            finally:
                # Clean up temp file
                if os.path.exists(tmp_path):
                    os.remove(tmp_path)
            
            # Ensure wav is numpy array
            if not isinstance(wav, np.ndarray):
                wav = np.array(wav)
            
            # Convert to mono if stereo
            if len(wav.shape) > 1:
                wav = np.mean(wav, axis=1)
            
            audio_segments.append({
                'audio': wav,
                'start_time': seg['start_time'],
                'duration': seg['duration']
            })
        except Exception as e:
            print(f"Warning: Failed to synthesize segment {i+1}: {e}")
            # Create silence for failed segments
            silence_duration = int(seg['duration'] * sample_rate)
            audio_segments.append({
                'audio': np.zeros(silence_duration, dtype=np.float32),
                'start_time': seg['start_time'],
                'duration': seg['duration']
            })
    
    # Get sample rate from TTS model
    try:
        if hasattr(tts, 'synthesizer') and hasattr(tts.synthesizer, 'output_sample_rate'):
            sample_rate = tts.synthesizer.output_sample_rate
        elif hasattr(tts, 'output_sample_rate'):
            sample_rate = tts.output_sample_rate
        else:
            sample_rate = 22050  # Default fallback
    except:
        sample_rate = 22050  # Default fallback
    
    # Combine all segments into a single audio file
    # Calculate total duration
    total_duration = max(seg['start_time'] + seg['duration'] for seg in audio_segments)
    total_samples = int(total_duration * sample_rate)
    combined_audio = np.zeros(total_samples, dtype=np.float32)
    
    # Place each segment at its original timestamp
    for seg_data in audio_segments:
        start_sample = int(seg_data['start_time'] * sample_rate)
        audio_data = seg_data['audio']
        
        # Ensure audio is 1D array
        if len(audio_data.shape) > 1:
            audio_data = audio_data.flatten()
        
        # Trim or pad to match expected duration
        expected_samples = int(seg_data['duration'] * sample_rate)
        if len(audio_data) > expected_samples:
            audio_data = audio_data[:expected_samples]
        elif len(audio_data) < expected_samples:
            # Pad with silence
            padding = np.zeros(expected_samples - len(audio_data), dtype=np.float32)
            audio_data = np.concatenate([audio_data, padding])
        
        # Place in combined audio
        end_sample = start_sample + len(audio_data)
        if end_sample <= len(combined_audio):
            combined_audio[start_sample:end_sample] = audio_data
    
    # Save audio file
    output_ext = os.path.splitext(output_path)[1].lower()
    if output_ext == '.wav':
        sf.write(output_path, combined_audio, sample_rate)
    elif output_ext == '.mp3':
        # Save as WAV first, then convert (requires ffmpeg)
        temp_wav = output_path.replace('.mp3', '.wav')
        sf.write(temp_wav, combined_audio, sample_rate)
        # Convert to MP3 using ffmpeg if available
        try:
            import subprocess
            subprocess.run(
                ['ffmpeg', '-i', temp_wav, '-codec:a', 'libmp3lame', '-qscale:a', '2', output_path],
                check=True,
                capture_output=True
            )
            os.remove(temp_wav)
        except:
            # If ffmpeg not available, keep WAV file
            print("Warning: ffmpeg not available. Saved as WAV instead of MP3.")
            output_path = temp_wav
    else:
        # Default to WAV
        if output_ext:
            output_path = os.path.splitext(output_path)[0] + '.wav'
        sf.write(output_path, combined_audio, sample_rate)
    
    print(f"Generated translated audio: {output_path}")
    return output_path


def _format_srt_timestamp(seconds: float) -> str:
    """Convert seconds to SRT timestamp format (HH:MM:SS,mmm)."""
    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    secs = int(seconds % 60)
    millis = int((seconds % 1) * 1000)
    return f"{hours:02d}:{minutes:02d}:{secs:02d},{millis:03d}"

